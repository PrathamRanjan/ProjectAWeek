# groq_smolvlm_analysis.py
import requests
import json
import base64
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.chrome.service import Service
from webdriver_manager.chrome import ChromeDriverManager
import time
import os
from PIL import Image

class GROQSmolVLMAnalyzer:
    def __init__(self, groq_api_key=None):
        self.groq_api_key = groq_api_key or "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"
        self.driver = None
        
    def setup_driver(self, headless=True):
        """Setup Chrome driver for screenshots"""
        print("🌐 Setting up Chrome driver...")
        
        chrome_options = Options()
        if headless:
            chrome_options.add_argument('--headless')
        
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument('--window-size=1920,1080')
        chrome_options.add_argument('--force-device-scale-factor=1')
        chrome_options.add_argument('--hide-scrollbars')
        
        service = Service(ChromeDriverManager().install())
        self.driver = webdriver.Chrome(service=service, options=chrome_options)
        print("✅ Chrome driver ready!")
    
    def take_screenshot(self, dashboard_url, output_filename=None, wait_time=15):
        """Take screenshot of the dashboard"""
        print("📸 Taking screenshot of PowerBI dashboard...")
        
        if not output_filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"dashboard_{timestamp}.png"
        
        try:
            self.driver.get(dashboard_url)
            print(f"⏳ Waiting {wait_time} seconds for dashboard to load...")
            time.sleep(wait_time)
            
            # Take screenshot
            self.driver.save_screenshot(output_filename)
            
            if os.path.exists(output_filename):
                file_size = os.path.getsize(output_filename) / 1024  # KB
                print(f"✅ Screenshot saved: {output_filename} ({file_size:.1f} KB)")
                return output_filename
            
        except Exception as e:
            print(f"❌ Screenshot failed: {e}")
        
        return None
    
    def encode_image_for_groq(self, image_path):
        """Encode image to base64 for GROQ API"""
        print("🔄 Encoding image for GROQ...")
        
        try:
            with open(image_path, "rb") as image_file:
                encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
            print("✅ Image encoded successfully")
            return encoded_image
        except Exception as e:
            print(f"❌ Failed to encode image: {e}")
            return None
    
    def extract_dashboard_info_with_smolvlm(self, image_path):
        """Use GROQ API with SmolVLM to extract dashboard information"""
        print("👁️ Analyzing dashboard with GROQ SmolVLM...")
        
        encoded_image = self.encode_image_for_groq(image_path)
        if not encoded_image:
            return None
        
        # Detailed prompt for SmolVLM vision analysis
        vlm_prompt = """Analyze this PowerBI dashboard image in extreme detail. Extract all visible information including:

1. TITLE: The main dashboard title/heading

2. KEY_METRICS: Every number, KPI, and metric visible with their exact values and labels
   - Example: "Total Vehicles: 95.50K", "Average Range: 33.71 kms"

3. CHARTS_ANALYSIS: For each chart/visualization:
   - Chart type (bar, line, pie, treemap, etc.)
   - Title of the chart
   - Data categories and values
   - Specific numbers visible
   - Any trends or patterns

4. FILTERS_CONTROLS: All filter controls, dropdowns, slicers visible

5. DATA_TABLES: Any tables with their exact content

6. GEOGRAPHIC_DATA: Any city/region names and their associated values

7. CATEGORIES: All dimensions shown (manufacturers, models, years, etc.)

8. PERCENTAGES: All percentage values visible

9. COMPARISONS: Any comparative data or breakdowns

10. SPECIFIC_VALUES: Every single number, percentage, count visible

Be extremely precise and extract every visible piece of data, text, and numbers from the dashboard."""
        
        # GROQ API call for SmolVLM
        groq_url = "https://api.groq.com/openai/v1/chat/completions"
        
        headers = {
            "Authorization": f"Bearer {self.groq_api_key}",
            "Content-Type": "application/json"
        }
        
        # Using GROQ's available vision model (Llama 4 Scout)
        payload = {
            "model": "meta-llama/llama-4-scout-17b-16e-instruct",  # GROQ vision model
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": vlm_prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{encoded_image}"
                            }
                        }
                    ]
                }
            ],
            "max_tokens": 3000,
            "temperature": 0.1
        }
        
        try:
            response = requests.post(groq_url, headers=headers, json=payload)
            response.raise_for_status()
            
            result = response.json()
            extracted_info = result['choices'][0]['message']['content']
            
            print("✅ GROQ SmolVLM analysis completed!")
            return extracted_info
            
        except Exception as e:
            print(f"❌ GROQ SmolVLM analysis failed: {e}")
            print(f"Response: {response.text if 'response' in locals() else 'No response'}")
            return None
    
    def analyze_with_groq_llama3(self, extracted_info):
        """Use GROQ Llama 3 to analyze the extracted dashboard information"""
        print("🤖 Analyzing with GROQ Llama 3...")
        
        analysis_prompt = f"""Based on the following detailed dashboard analysis from SmolVLM, provide a comprehensive business analysis:

DASHBOARD_INFORMATION:
{extracted_info}

Structure your response with these exact section headers:

**EXECUTIVE_SUMMARY**
Provide a brief overview of the dashboard purpose and key findings.

**KEY_METRICS**
Analyze the main KPIs, their performance, and business significance. Reference specific numbers.

**VISUAL_ANALYSIS**
Interpret each chart/visualization, explaining what they reveal about business performance.

**TRENDS_PATTERNS**
Identify trends, patterns, and notable insights from the data.

**INSIGHTS**
Provide actionable business insights based on the actual data.

**RECOMMENDATIONS**
Suggest specific, actionable recommendations based on the dashboard findings.

Reference specific numbers, percentages, and values from the vision analysis whenever possible."""
        
        # GROQ API call for Llama 3
        groq_url = "https://api.groq.com/openai/v1/chat/completions"
        
        headers = {
            "Authorization": f"Bearer {self.groq_api_key}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": "llama3-8b-8192",
            "messages": [
                {
                    "role": "system",
                    "content": "You are an expert business analyst. Analyze dashboard data and provide actionable insights with specific references to the numbers and metrics provided."
                },
                {
                    "role": "user",
                    "content": analysis_prompt
                }
            ],
            "max_tokens": 2500,
            "temperature": 0.3
        }
        
        try:
            response = requests.post(groq_url, headers=headers, json=payload)
            response.raise_for_status()
            
            result = response.json()
            analysis = result['choices'][0]['message']['content']
            
            print("✅ GROQ Llama 3 analysis completed!")
            return analysis
            
        except Exception as e:
            print(f"❌ GROQ Llama 3 analysis failed: {e}")
            return None
    
    def extract_sections_from_analysis(self, analysis_text):
        """Extract structured content from analysis for slides"""
        print("📝 Extracting slide content...")
        
        if not analysis_text:
            return {}
        
        sections = {
            'executive_summary': '',
            'key_metrics': '',
            'visual_analysis': '',
            'trends_patterns': '',
            'insights': '',
            'recommendations': ''
        }
        
        lines = analysis_text.split('\n')
        current_section = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Check for section headers
            if '**EXECUTIVE_SUMMARY**' in line:
                current_section = 'executive_summary'
                continue
            elif '**KEY_METRICS**' in line:
                current_section = 'key_metrics'
                continue
            elif '**VISUAL_ANALYSIS**' in line:
                current_section = 'visual_analysis'
                continue
            elif '**TRENDS_PATTERNS**' in line:
                current_section = 'trends_patterns'
                continue
            elif '**INSIGHTS**' in line:
                current_section = 'insights'
                continue
            elif '**RECOMMENDATIONS**' in line:
                current_section = 'recommendations'
                continue
            elif current_section and line and not line.startswith('**'):
                sections[current_section] += line + '\n'
        
        # Clean up sections
        for key in sections:
            sections[key] = sections[key].strip()
        
        return sections
    
    def optimize_image_for_powerpoint(self, image_path, max_width=7, max_height=5):
        """Optimize image for PowerPoint"""
        print("🖼️ Optimizing image for PowerPoint...")
        
        try:
            with Image.open(image_path) as img:
                width, height = img.size
                aspect_ratio = width / height
                
                max_width_px = int(max_width * 96)
                max_height_px = int(max_height * 96)
                
                if width > max_width_px or height > max_height_px:
                    if aspect_ratio > max_width / max_height:
                        new_width = max_width_px
                        new_height = int(new_width / aspect_ratio)
                    else:
                        new_height = max_height_px
                        new_width = int(new_height * aspect_ratio)
                    
                    img_resized = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
                    optimized_path = image_path.replace('.png', '_optimized.png')
                    img_resized.save(optimized_path, 'PNG', quality=95)
                    
                    print(f"✅ Image optimized: {new_width}x{new_height}")
                    return optimized_path
                else:
                    print("✅ Image already optimal size")
                    return image_path
                    
        except Exception as e:
            print(f"❌ Image optimization failed: {e}")
            return image_path
    
    def create_powerpoint_from_analysis(self, image_path, analysis_sections, extracted_info=None, title="Dashboard Analysis Report"):
        """Create PowerPoint presentation"""
        print("📄 Creating PowerPoint presentation...")
        
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ppt_filename = f"dashboard_analysis_{timestamp}.pptx"
        
        try:
            prs = Presentation()
            
            # Title Slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title
            subtitle = title_slide.placeholders[1]
            subtitle.text = f"GROQ SmolVLM + Llama 3 Analysis\nGenerated: {datetime.now().strftime('%B %d, %Y')}"
            
            # Executive Summary
            exec_slide = prs.slides.add_slide(prs.slide_layouts[1])
            exec_slide.shapes.title.text = "Executive Summary"
            exec_slide.placeholders[1].text = analysis_sections.get('executive_summary', 'Analysis in progress...')
            
            # Dashboard Overview with Image
            dashboard_slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Add title
            title_shape = dashboard_slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            title_frame = title_shape.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = "Dashboard Overview"
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            title_para.font.color.rgb = RGBColor(0, 51, 102)
            
            # Add screenshot
            if os.path.exists(image_path):
                optimized_image = self.optimize_image_for_powerpoint(image_path)
                if os.path.exists(optimized_image):
                    dashboard_slide.shapes.add_picture(
                        optimized_image, 
                        Inches(1), 
                        Inches(1.2), 
                        width=Inches(8)
                    )
                    print("✅ Dashboard image added to slide")
            
            # Content slides
            slides_content = [
                ("Key Metrics Analysis", analysis_sections.get('key_metrics', '')),
                ("Visual Analysis", analysis_sections.get('visual_analysis', '')),
                ("Trends & Patterns", analysis_sections.get('trends_patterns', '')),
                ("Business Insights", analysis_sections.get('insights', '')),
                ("Recommendations", analysis_sections.get('recommendations', ''))
            ]
            
            for slide_title, slide_content in slides_content:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_title
                slide.placeholders[1].text = slide_content or 'Content in progress...'
            
            # SmolVLM Raw Analysis (Optional)
            if extracted_info:
                data_slide = prs.slides.add_slide(prs.slide_layouts[1])
                data_slide.shapes.title.text = "SmolVLM Raw Analysis"
                truncated_info = extracted_info[:1500] + "..." if len(extracted_info) > 1500 else extracted_info
                data_slide.placeholders[1].text = truncated_info
            
            # Format slides
            for i in range(1, len(prs.slides)):
                slide = prs.slides[i]
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.text:
                                paragraph.font.size = Pt(14)
                                paragraph.font.name = 'Calibri'
                                paragraph.space_after = Pt(6)
            
            # Save presentation
            prs.save(ppt_filename)
            
            if os.path.exists(ppt_filename):
                file_size = os.path.getsize(ppt_filename) / 1024
                print(f"✅ PowerPoint created: {ppt_filename} ({file_size:.1f} KB)")
                return ppt_filename
                
        except Exception as e:
            print(f"❌ PowerPoint creation failed: {e}")
            return None
        
        return None
    
    def close(self):
        """Clean up resources"""
        if self.driver:
            self.driver.quit()
            print("🔒 Browser closed")

def create_groq_smolvlm_report(dashboard_url, title="Dashboard Analysis Report"):
    """
    Complete workflow: Screenshot → GROQ SmolVLM → GROQ Llama 3 → PowerPoint
    
    Args:
        dashboard_url: PowerBI dashboard URL
        title: Presentation title
    
    Returns:
        tuple: (screenshot_path, extracted_info, analysis_text, powerpoint_path)
    """
    
    analyzer = GROQSmolVLMAnalyzer()
    
    try:
        print("🚀 GROQ SmolVLM + Llama 3 Analysis Pipeline")
        print("=" * 60)
        
        # Step 1: Take screenshot
        analyzer.setup_driver(headless=True)
        screenshot_path = analyzer.take_screenshot(dashboard_url)
        
        if not screenshot_path:
            print("❌ Failed to capture screenshot")
            return None, None, None, None
        
        # Step 2: Analyze with GROQ SmolVLM
        extracted_info = analyzer.extract_dashboard_info_with_smolvlm(screenshot_path)
        
        if not extracted_info:
            print("❌ GROQ SmolVLM analysis failed")
            return screenshot_path, None, None, None
        
        # Step 3: Analyze with GROQ Llama 3
        analysis_text = analyzer.analyze_with_groq_llama3(extracted_info)
        
        if not analysis_text:
            print("❌ GROQ Llama 3 analysis failed")
            return screenshot_path, extracted_info, None, None
        
        # Step 4: Extract sections and create PowerPoint
        analysis_sections = analyzer.extract_sections_from_analysis(analysis_text)
        powerpoint_path = analyzer.create_powerpoint_from_analysis(
            screenshot_path, 
            analysis_sections,
            extracted_info,
            title
        )
        
        return screenshot_path, extracted_info, analysis_text, powerpoint_path
        
    finally:
        analyzer.close()

# Main execution
if __name__ == "__main__":
    print("🚀 GROQ SmolVLM + Llama 3 Dashboard Analyzer")
    print("=" * 60)
    
    # Dashboard URL
    dashboard_url = "https://app.fabric.microsoft.com/view?r=eyJrIjoiODhmOThlYWMtYmI0MS00MmIyLWE4YjktYmIzM2EwYjAzYTBiIiwidCI6ImVlMmQ2ZDcyLTk1MzUtNDI0Mi1hMDc3LWFjZjE4NTc4MmY5YiIsImMiOjF9"
    
    # Create report
    screenshot_path, extracted_info, analysis, powerpoint_path = create_groq_smolvlm_report(
        dashboard_url=dashboard_url,
        title="Electric Vehicle Dashboard Analysis"
    )
    
    if screenshot_path and extracted_info and analysis and powerpoint_path:
        print(f"\n🎉 SUCCESS! GROQ SmolVLM + Llama 3 report created!")
        print(f"📸 Screenshot: {screenshot_path}")
        print(f"👁️ SmolVLM Analysis: {len(extracted_info)} characters")
        print(f"🤖 Llama 3 Analysis: {len(analysis)} characters")
        print(f"📄 PowerPoint: {powerpoint_path}")
        
        print(f"\n📋 Pipeline completed:")
        print("✅ 1. Screenshot captured")
        print("✅ 2. GROQ SmolVLM vision analysis")
        print("✅ 3. GROQ Llama 3 business analysis")
        print("✅ 4. PowerPoint generated")
        
        print(f"\n💡 SmolVLM Analysis Preview:")
        print(extracted_info[:300] + "..." if len(extracted_info) > 300 else extracted_info)
        
        print(f"\n🏆 Ready for competition!")
    else:
        print("\n❌ Pipeline failed")
        if screenshot_path:
            print(f"✅ Screenshot: {screenshot_path}")
        if extracted_info:
            print(f"✅ SmolVLM extracted: {len(extracted_info)} chars")
        if analysis:
            print(f"✅ Llama 3 analyzed: {len(analysis)} chars")